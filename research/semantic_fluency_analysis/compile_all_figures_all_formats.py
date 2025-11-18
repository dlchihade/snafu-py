#!/usr/bin/env python3
"""Compile ALL figures in ALL formats (.svg, .png, .pdf) into a PowerPoint presentation"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os

def find_all_figure_files():
    """Find all figure files in all formats"""
    output_dir = Path('output')
    
    # Base figure names (without extension)
    base_figures = [
        # Main NATURE_REAL figures
        'NATURE_REAL_figure0_alpha_violin_swarm',
        'NATURE_REAL_figure0b_coherence_violin_swarm',
        'NATURE_REAL_figure1_exploration_exploitation',
        'NATURE_REAL_figure2_phase_coherence',
        'NATURE_REAL_figure3_meg_correlations',
        'NATURE_REAL_figure4_comprehensive',
        'NATURE_REAL_behavior_performance',
        'NATURE_REAL_ee_vs_moca',
        'NATURE_REAL_compiled_results',
        'NATURE_REAL_participant_values_table',
        'NATURE_REAL_violin_lc_residuals',
        
        # Intermediate figures
        'figures/intermediate/participant_demographics',
        'figures/intermediate/age_distribution',
        'figures/intermediate/alpha_power_violin',
        'figures/intermediate/lc_vs_alpha',
        'figures/intermediate/svf_vs_ee',
        'figures/intermediate/svf_ee_boxswarm',
        
        # Additional analysis figures
        'alpha_vs_LC_residuals',
        'semantic_fluency_analysis',
        'compiled_distribution_figures',
        'exploit_explore_bar',
        'combined_mediation_exploit_panels',
        'test_grid_panels',
        
        # Mediation figures
        'mediation_svf_age_nature',
        'mediation_exploit_age_nature',
        'mediation_exploit_coherence_metric_nature',
        'mediation_svf_disease_stage_working',
        'mediation_ee_disease_stage_working',
        'mediation_age_vs_disease',
    ]
    
    # Also check for SVG files in figures directory
    svg_figures = [
        'figures/schematic_diagram',
        'figures/NATURE_REAL_figure4_panels',
        'distribution_svf_scores',
        'distribution_total_words',
    ]
    
    all_figures = []
    
    # Find all formats for each base figure
    for base_name in base_figures + svg_figures:
        base_path = output_dir / base_name if '/' not in base_name else output_dir / base_name
        
        # Check for all formats
        formats_found = []
        for ext in ['.png', '.pdf', '.svg']:
            file_path = base_path.parent / f"{base_path.name}{ext}"
            if file_path.exists():
                formats_found.append((ext, file_path))
        
        if formats_found:
            # Create a readable title from the base name
            title = base_name.replace('NATURE_REAL_', '').replace('_', ' ').title()
            if 'figures/' in base_name:
                title = 'Intermediate: ' + base_name.split('/')[-1].replace('_', ' ').title()
            
            all_figures.append((title, formats_found))
    
    return all_figures, output_dir

def create_presentation():
    """Create PowerPoint with all figures in all formats"""
    all_figures, output_dir = find_all_figure_files()
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Semantic Fluency Analysis"
    subtitle.text = "All Figures - All Formats\n(SVG, PNG, PDF)\nExcluding Zero Values (N=45)"
    
    total_slides = 0
    
    # Add each figure format as a slide
    for fig_title, formats_found in all_figures:
        for ext, fig_path in formats_found:
            # Skip if file doesn't exist
            if not fig_path.exists():
                continue
            
            # Create blank slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            total_slides += 1
            
            # Try to get image dimensions (works for PNG, may not work for PDF/SVG)
            try:
                if ext == '.png':
                    img = Image.open(fig_path)
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height
                elif ext == '.pdf':
                    # For PDF, use default aspect ratio or try to read
                    aspect_ratio = 1.4  # Default for most figures
                elif ext == '.svg':
                    # For SVG, use default aspect ratio
                    aspect_ratio = 1.4
                else:
                    aspect_ratio = 1.4
            except Exception:
                aspect_ratio = 1.4
            
            # Calculate size to fit slide (with margins)
            slide_width = prs.slide_width - Inches(0.5)
            slide_height = prs.slide_height - Inches(1.0)
            
            if aspect_ratio > (slide_width / slide_height):
                # Image is wider - fit to width
                width = slide_width
                height = slide_width / aspect_ratio
            else:
                # Image is taller - fit to height
                height = slide_height
                width = slide_height * aspect_ratio
            
            # Center the image
            left = (prs.slide_width - width) / 2
            top = Inches(0.5)
            
            # Add image (for PDF/SVG, this may not work perfectly, but we'll try)
            try:
                if ext == '.png':
                    slide.shapes.add_picture(str(fig_path), left, top, width, height)
                elif ext == '.pdf':
                    # PDFs can't be directly added, but we can note it
                    # For now, skip PDFs or convert them
                    print(f"⚠️  PDF format not directly supported: {fig_path.name}")
                    continue
                elif ext == '.svg':
                    # SVGs also can't be directly added
                    print(f"⚠️  SVG format not directly supported: {fig_path.name}")
                    continue
            except Exception as e:
                print(f"⚠️  Could not add {fig_path.name}: {e}")
                continue
            
            # Add title and file path as text box at bottom
            left_text = Inches(0.5)
            top_text = prs.slide_height - Inches(0.5)
            width_text = prs.slide_width - Inches(1.0)
            height_text = Inches(0.4)
            
            textbox = slide.shapes.add_textbox(left_text, top_text, width_text, height_text)
            text_frame = textbox.text_frame
            text_frame.text = f"{fig_title} ({ext.upper()})"
            text_frame.paragraphs[0].font.size = Inches(0.12)
            text_frame.paragraphs[0].font.bold = True
            
            # Add file path in smaller text
            p = text_frame.add_paragraph()
            p.text = str(fig_path.relative_to(output_dir.parent))
            p.font.size = Inches(0.08)
            p.font.italic = True
            
            print(f"✅ Added: {fig_title} ({ext})")
    
    # Save presentation
    output_path = output_dir / 'all_figures_all_formats_presentation.pptx'
    prs.save(str(output_path))
    print(f"\n🎉 PowerPoint saved: {output_path}")
    print(f"   Total slides: {total_slides + 1} (including title slide)")
    return output_path

if __name__ == '__main__':
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx not installed")
        print("   Install with: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()

