#!/usr/bin/env python3
"""Compile all figures into a PowerPoint presentation"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os

def get_figure_files():
    """Get all figure files in order"""
    output_dir = Path('output')
    
    figures = [
        # Main NATURE_REAL figures
        ('Figure 0a: Alpha Power Distribution', 'NATURE_REAL_figure0_alpha_violin_swarm.png'),
        ('Figure 0b: Coherence Distribution', 'NATURE_REAL_figure0b_coherence_violin_swarm.png'),
        ('Figure 1: Exploration vs Exploitation', 'NATURE_REAL_figure1_exploration_exploitation.png'),
        ('Figure 2: Phase Coherence', 'NATURE_REAL_figure2_phase_coherence.png'),
        ('Figure 3: MEG Correlations', 'NATURE_REAL_figure3_meg_correlations.png'),
        ('Figure 4: Comprehensive Scatter Plots', 'NATURE_REAL_figure4_comprehensive.png'),
        ('Behavior Performance Metrics', 'NATURE_REAL_behavior_performance.png'),
        ('E-E Index vs MoCA', 'NATURE_REAL_ee_vs_moca.png'),
        ('Compiled Results', 'NATURE_REAL_compiled_results.png'),
        ('Participant Demographics Table', 'NATURE_REAL_participant_values_table.png'),
        ('LC Residuals Violin Plot', 'NATURE_REAL_violin_lc_residuals.png'),
        
        # Intermediate figures
        ('Intermediate: Participant Demographics', 'figures/intermediate/participant_demographics.png'),
        ('Intermediate: Age Distribution', 'figures/intermediate/age_distribution.png'),
        ('Intermediate: Alpha Power Violin', 'figures/intermediate/alpha_power_violin.png'),
        ('Intermediate: LC vs Alpha', 'figures/intermediate/lc_vs_alpha.png'),
        ('Intermediate: SVF vs EE', 'figures/intermediate/svf_vs_ee.png'),
        ('Intermediate: SVF EE Box-Swarm', 'figures/intermediate/svf_ee_boxswarm.png'),
        
        # Additional analysis figures
        ('Alpha vs LC Residuals', 'alpha_vs_LC_residuals.png'),
        ('Semantic Fluency Analysis', 'semantic_fluency_analysis.png'),
        ('Compiled Distribution Figures', 'compiled_distribution_figures.png'),
        ('Exploit Explore Bar Chart', 'exploit_explore_bar.png'),
        ('Combined Mediation Exploit Panels', 'combined_mediation_exploit_panels.png'),
        ('Test Grid Panels', 'test_grid_panels.png'),
        
        # Mediation figures
        ('Mediation: SVF (Age-adjusted)', 'mediation_svf_age_nature.png'),
        ('Mediation: EE Metric (Age-adjusted)', 'mediation_exploit_age_nature.png'),
        ('Mediation: EE Coherence Metric (Age-adjusted)', 'mediation_exploit_coherence_metric_nature.png'),
        ('Mediation: SVF (Disease Stage)', 'mediation_svf_disease_stage_working.png'),
        ('Mediation: EE (Disease Stage)', 'mediation_ee_disease_stage_working.png'),
        ('Mediation: Age vs Disease', 'mediation_age_vs_disease.png'),
    ]
    
    return figures, output_dir

def create_presentation():
    """Create PowerPoint with all figures"""
    figures, output_dir = get_figure_files()
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Semantic Fluency Analysis"
    subtitle.text = "All Figures\nExcluding Zero Values (N=45)"
    
    # Add each figure as a slide
    for fig_title, fig_filename in figures:
        # Handle both direct paths and subdirectory paths
        if '/' in fig_filename:
            fig_path = output_dir / fig_filename
        else:
            fig_path = output_dir / fig_filename
        
        if not fig_path.exists():
            print(f"⚠️  Skipping {fig_filename} (not found at {fig_path})")
            continue
        
        # Create blank slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Get image dimensions
        img = Image.open(fig_path)
        img_width, img_height = img.size
        aspect_ratio = img_width / img_height
        
        # Calculate size to fit slide (with margins)
        slide_width = prs.slide_width - Inches(0.5)
        slide_height = prs.slide_height - Inches(1.0)
        
        if aspect_ratio > (slide_width / slide_height):
            # Image is wider - fit to width
            width = slide_width
            height = slide_width / aspect_ratio
        else:
            # Image is taller - fit to height
            height = slide_height
            width = slide_height * aspect_ratio
        
        # Center the image
        left = (prs.slide_width - width) / 2
        top = Inches(0.5)
        
        # Add image
        slide.shapes.add_picture(str(fig_path), left, top, width, height)
        
        # Add title as text box at bottom
        left_text = Inches(0.5)
        top_text = prs.slide_height - Inches(0.4)
        width_text = prs.slide_width - Inches(1.0)
        height_text = Inches(0.3)
        
        textbox = slide.shapes.add_textbox(left_text, top_text, width_text, height_text)
        text_frame = textbox.text_frame
        text_frame.text = fig_title
        text_frame.paragraphs[0].font.size = Inches(0.15)
        text_frame.paragraphs[0].font.bold = True
        
        print(f"✅ Added: {fig_title}")
    
    # Save presentation
    output_path = output_dir / 'all_figures_presentation.pptx'
    prs.save(str(output_path))
    print(f"\n🎉 PowerPoint saved: {output_path}")
    return output_path

if __name__ == '__main__':
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx not installed")
        print("   Install with: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()

