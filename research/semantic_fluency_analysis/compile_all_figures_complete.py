#!/usr/bin/env python3
"""Compile ALL figures in ALL formats (.svg, .png, .pdf) into a PowerPoint with file paths"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os

def find_all_figure_files():
    """Find all figure files in all formats"""
    output_dir = Path('output')
    
    # Base figure names (without extension)
    base_figures = [
        # Main NATURE_REAL figures
        'NATURE_REAL_figure0_alpha_violin_swarm',
        'NATURE_REAL_figure0b_coherence_violin_swarm',
        'NATURE_REAL_figure1_exploration_exploitation',
        'NATURE_REAL_figure2_phase_coherence',
        'NATURE_REAL_figure3_meg_correlations',
        'NATURE_REAL_figure4_comprehensive',
        'NATURE_REAL_behavior_performance',
        'NATURE_REAL_ee_vs_moca',
        'NATURE_REAL_compiled_results',
        'NATURE_REAL_participant_values_table',
        'NATURE_REAL_violin_lc_residuals',
        
        # Intermediate figures
        'figures/intermediate/participant_demographics',
        'figures/intermediate/age_distribution',
        'figures/intermediate/alpha_power_violin',
        'figures/intermediate/lc_vs_alpha',
        'figures/intermediate/svf_vs_ee',
        'figures/intermediate/svf_ee_boxswarm',
        
        # Additional analysis figures
        'alpha_vs_LC_residuals',
        'semantic_fluency_analysis',
        'compiled_distribution_figures',
        'exploit_explore_bar',
        'combined_mediation_exploit_panels',
        'test_grid_panels',
        
        # Mediation figures
        'mediation_svf_age_nature',
        'mediation_exploit_age_nature',
        'mediation_exploit_coherence_metric_nature',
        'mediation_svf_disease_stage_working',
        'mediation_ee_disease_stage_working',
        'mediation_age_vs_disease',
    ]
    
    # Also check for SVG files
    svg_figures = [
        'figures/schematic_diagram',
        'figures/NATURE_REAL_figure4_panels',
        'distribution_svf_scores',
        'distribution_total_words',
    ]
    
    all_figures = []
    
    # Find all formats for each base figure
    for base_name in base_figures + svg_figures:
        base_path = output_dir / base_name if '/' not in base_name else output_dir / base_name
        
        # Check for all formats
        formats_found = []
        for ext in ['.png', '.pdf', '.svg']:
            file_path = base_path.parent / f"{base_path.name}{ext}"
            if file_path.exists():
                formats_found.append((ext, file_path))
        
        if formats_found:
            # Create a readable title from the base name
            title = base_name.replace('NATURE_REAL_', '').replace('_', ' ').title()
            if 'figures/' in base_name:
                title = 'Intermediate: ' + base_name.split('/')[-1].replace('_', ' ').title()
            
            all_figures.append((title, formats_found))
    
    return all_figures, output_dir

def create_presentation():
    """Create PowerPoint with all figures in all formats, including file paths"""
    all_figures, output_dir = find_all_figure_files()
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Semantic Fluency Analysis"
    subtitle.text = "All Figures - All Formats\n(SVG, PNG, PDF with File Paths)\nExcluding Zero Values (N=45)"
    
    total_slides = 0
    
    # Add each figure format as a slide
    for fig_title, formats_found in all_figures:
        # Sort formats: PNG first (best for display), then PDF, then SVG
        formats_sorted = sorted(formats_found, key=lambda x: (x[0] != '.png', x[0] != '.pdf', x[0]))
        
        for ext, fig_path in formats_sorted:
            # Skip if file doesn't exist
            if not fig_path.exists():
                continue
            
            # Create blank slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            total_slides += 1
            
            # Try to add image (only works for PNG)
            image_added = False
            if ext == '.png':
                try:
                    img = Image.open(fig_path)
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height
                    
                    # Calculate size to fit slide (with margins)
                    slide_width = prs.slide_width - Inches(0.5)
                    slide_height = prs.slide_height - Inches(1.2)  # More space for text
                    
                    if aspect_ratio > (slide_width / slide_height):
                        width = slide_width
                        height = slide_width / aspect_ratio
                    else:
                        height = slide_height
                        width = slide_height * aspect_ratio
                    
                    # Center the image
                    left = (prs.slide_width - width) / 2
                    top = Inches(0.5)
                    
                    slide.shapes.add_picture(str(fig_path), left, top, width, height)
                    image_added = True
                except Exception as e:
                    print(f"⚠️  Could not add PNG {fig_path.name}: {e}")
            
            # If no image was added (PDF/SVG), add a placeholder text box
            if not image_added:
                left_box = Inches(1)
                top_box = Inches(2)
                width_box = prs.slide_width - Inches(2)
                height_box = Inches(3)
                
                textbox = slide.shapes.add_textbox(left_box, top_box, width_box, height_box)
                text_frame = textbox.text_frame
                text_frame.text = f"{ext.upper()} Format\n\nThis format cannot be directly displayed in PowerPoint.\nPlease open the file using the path below."
                text_frame.paragraphs[0].font.size = Inches(0.15)
                text_frame.paragraphs[0].font.bold = True
                text_frame.paragraphs[0].alignment = 1  # Center alignment
            
            # Add title and file path as text box at bottom
            left_text = Inches(0.5)
            top_text = prs.slide_height - Inches(0.6)
            width_text = prs.slide_width - Inches(1.0)
            height_text = Inches(0.5)
            
            textbox = slide.shapes.add_textbox(left_text, top_text, width_text, height_text)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Title
            p = text_frame.paragraphs[0]
            p.text = f"{fig_title} ({ext.upper()})"
            p.font.size = Inches(0.12)
            p.font.bold = True
            
            # File path
            p2 = text_frame.add_paragraph()
            p2.text = f"File: {fig_path.relative_to(output_dir.parent)}"
            p2.font.size = Inches(0.08)
            p2.font.italic = True
            
            # Absolute path
            p3 = text_frame.add_paragraph()
            p3.text = f"Full path: {fig_path.absolute()}"
            p3.font.size = Inches(0.07)
            p3.font.name = 'Courier New'
            
            print(f"✅ Added: {fig_title} ({ext}) - {fig_path.relative_to(output_dir.parent)}")
    
    # Save presentation
    output_path = output_dir / 'all_figures_all_formats_complete.pptx'
    prs.save(str(output_path))
    print(f"\n🎉 PowerPoint saved: {output_path}")
    print(f"   Total slides: {total_slides + 1} (including title slide)")
    print(f"   Includes: PNG (displayed), PDF (paths only), SVG (paths only)")
    return output_path

if __name__ == '__main__':
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx not installed")
        print("   Install with: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()

